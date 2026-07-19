import os


def _norm(s: str) -> str:
    """Lowercase + fold ё→е so keyword matching is accent-insensitive."""
    return s.lower().replace("ё", "е")


# Fallback used when the organizer has not configured any checklist items.
# Each entry is (keyword searched in slides, human label). The keyword is also
# the canonical key returned in found_topics/missing_topics (the frontend and
# seed data both match on the key, not the label).
DEFAULT_TOPICS = [
    ("проблема", "Проблема"),
    ("решение", "Решение"),
    ("аудитория", "Целевая аудитория"),
    ("стек", "Технологический стек"),
    ("демо", "Демо"),
    ("команда", "Команда"),
    ("контакты", "Контакты"),
]
MIN_SLIDES = 8


def _check_pptx(file_path: str) -> tuple[int, str]:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return len(prs.slides), " ".join(texts).lower()


def _check_pdf(file_path: str) -> tuple[int, str]:
    import PyPDF2
    text_parts = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")
    return len(reader.pages), " ".join(text_parts).lower()


def check_presentation(file_path: str, topics: list | None = None) -> dict:
    """`topics` is a list of (keyword, label) pairs from the organizer's
    checklist; falls back to DEFAULT_TOPICS when not provided."""
    topics = topics or DEFAULT_TOPICS
    ext = os.path.splitext(file_path)[-1].lower()
    try:
        if ext == ".pptx":
            slide_count, text = _check_pptx(file_path)
        else:
            slide_count, text = _check_pdf(file_path)
    except Exception as exc:
        return {"passed": False, "score": 0.0, "details": {"error": str(exc)}}

    text_norm = _norm(text)
    found_topics = [keyword for keyword, _label in topics if _norm(keyword) in text_norm]
    missing_topics = [keyword for keyword, _label in topics if _norm(keyword) not in text_norm]
    total_topics = len(topics) or 1

    slide_score = min(slide_count, 15) * 0.4            # max 6.0
    topic_score = (len(found_topics) / total_topics) * 4.0  # max 4.0
    score = min(slide_score + topic_score, 10.0)

    passed = slide_count >= MIN_SLIDES and len(found_topics) >= 3

    return {
        "passed": passed,
        "score": round(score, 1),
        "details": {
            "slide_count": slide_count,
            "found_topics": found_topics,
            "missing_topics": missing_topics,
            "min_slides": MIN_SLIDES,
        },
    }
